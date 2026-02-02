import streamlit as st
from google import genai
from google.genai import types
from PIL import Image
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import re

# ページ設定
st.set_page_config(page_title="PK-Insight Canvas", layout="wide")

# スライド作成関数
def create_pptx(strategy_text, client_name, product_name):
    prs = Presentation()
    # 表紙
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{product_name} 戦略提案書"
    slide.placeholders[1].text = f"Client: {client_name}\nCreated by PK-Insight Canvas"

    # 戦略セクションを分割してスライド化
    sections = re.split(r'\n(?=\d\.)', strategy_text)
    for section in sections:
        if not section.strip(): continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = section.strip().split('\n')
        slide.shapes.title.text = lines[0]
        slide.placeholders[1].text = "\n".join(lines[1:])

    ppt_io = BytesIO()
    prs.save(ppt_io)
    return ppt_io.getvalue()

# --- UI ---
st.title("🚀 PK-Insight Canvas v0.2 (Stable)")

with st.sidebar:
    st.header("⚙️ Settings")
    api_key = st.text_input("Gemini API Key", type="password")
    st.divider()
    client_name = st.text_input("クライアント名", "大手自動車メーカー")
    product_name = st.text_input("対象商品", "新型EV")
    target_user = st.text_area("ターゲット", "30代、都心、先進層")
    feedback = st.text_area("追加要望", "先進的な未来感と信頼性の両立")
    generate_btn = st.button("Generate Strategy & Slide")

if generate_btn:
    if not api_key:
        st.error("APIキーを入力してください")
    else:
        try:
            client = genai.Client(api_key=api_key)
            
            # ステップ1: 戦略テキストの生成（1.5 Flashを使用）
            with st.spinner("戦略ロジックを構築中..."):
                text_prompt = f"{client_name}の{product_name}に関する上申用戦略(1-8の項目)を作成してください。ターゲットは{target_user}、要望は{feedback}です。"
                # 無料枠の制限が緩い gemini-1.5-flash を指定
                text_response = client.models.generate_content(
                    model="gemini-1.5-flash", 
                    contents=text_prompt
                )
                strategy_text = text_response.text

            # 表示
            st.subheader("📊 Strategic Logic")
            st.write(strategy_text)

            # PPTX生成
            st.divider()
            pptx_data = create_pptx(strategy_text, client_name, product_name)
            st.download_button(
                label="📥 PowerPointをダウンロード", 
                data=pptx_data, 
                file_name=f"{product_name}_戦略案.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            st.success("スライドの書き出しが完了しました。")

        except Exception as e:
            st.error(f"エラーが発生しました: {e}")
            st.info("数十秒待ってから再度お試しいただくか、APIキーの制限を確認してください。")
